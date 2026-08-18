#!/usr/bin/env python3
"""Post-processing script for SlideSage-generated PPTX files.

Fixes fonts, applies custom styling, and enhances layout.

Usage:
    cd /tmp && source pptx-env/bin/activate
    python pptx-post-processing.py input.pptx output.pptx

Setup:
    cd /tmp && uv venv pptx-env && source pptx-env/bin/activate && uv pip install python-pptx
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color palette - minimalist luxury (white/beige)
CREAM = RGBColor(0xF8, 0xF5, 0xF0)
DARK = RGBColor(0x2C, 0x24, 0x1E)
BODY = RGBColor(0x4A, 0x45, 0x3E)
ACCENT = RGBColor(0x8B, 0x73, 0x55)

# Font settings
CN_FONT = "PingFang SC"  # Critical for Chinese content
EN_FONT = "Georgia"      # Serif for English titles

def fix_fonts(prs):
    """Fix fonts across all slides, especially for Chinese content."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = CN_FONT

def apply_color_scheme(prs, bg_color=CREAM):
    """Apply consistent color scheme."""
    for slide in prs.slides:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

def style_titles(prs):
    """Apply consistent title styling."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            
            is_title = len(text) < 60 and any(kw in text for kw in [
                'ANTOKEN', '赛道', '产品', '市场', '团队', '融资',
                '资金', '痛点', '方案', '数据', '模式', '里程碑'
            ])
            
            for para in shape.text_frame.paragraphs:
                if is_title:
                    para.font.size = Pt(22)
                    para.font.color.rgb = DARK
                    para.font.bold = True
                    para.font.name = CN_FONT
                elif text.startswith('•') or text.startswith('-'):
                    para.font.size = Pt(11)
                    para.font.color.rgb = BODY
                    para.font.name = CN_FONT
                    para.space_after = Pt(6)

def process(input_path, output_path):
    """Main processing function."""
    prs = Presentation(input_path)
    print(f"Processing {len(prs.slides)} slides...")
    
    apply_color_scheme(prs)
    fix_fonts(prs)
    style_titles(prs)
    
    prs.save(output_path)
    print(f"Saved to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pptx-post-processing.py input.pptx output.pptx")
        sys.exit(1)
    process(sys.argv[1], sys.argv[2])
